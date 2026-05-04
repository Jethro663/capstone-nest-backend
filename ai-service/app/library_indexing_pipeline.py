from __future__ import annotations

import hashlib
import logging
from pathlib import Path
from typing import Any

import fitz
from sqlalchemy import bindparam, text as sa_text
from sqlalchemy.dialects import postgresql
from sqlalchemy.ext.asyncio import AsyncSession

from . import ollama_client
from .backend_uploads import materialize_backend_upload
from .config import settings
from .content_sanitizer import sanitize_extracted_text
from .embedding_provider import embed_texts, embedding_to_vector_literal
from .indexing_pipeline import chunk_text_for_indexing, estimate_token_count

logger = logging.getLogger(__name__)


def _is_within(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
        return True
    except ValueError:
        return False


def _resolve_uploaded_path(file_path: str) -> Path:
    normalized = (file_path or "").strip()
    upload_root = Path(settings.upload_dir).resolve()
    normalized_slash = normalized.replace("\\", "/").lstrip("./")
    if normalized_slash.startswith("uploads/"):
        normalized_slash = normalized_slash[len("uploads/") :]

    candidates: list[Path] = []
    raw_candidate = Path(normalized)
    if raw_candidate.is_absolute():
        candidates.append(raw_candidate)

    candidates.extend(
        [
            raw_candidate,
            upload_root / normalized_slash,
            upload_root / Path(normalized).name,
        ]
    )

    deduped_candidates: list[Path] = []
    seen: set[Path] = set()
    for candidate in candidates:
        resolved = candidate.resolve()
        if resolved in seen:
            continue
        seen.add(resolved)
        if _is_within(resolved, upload_root):
            deduped_candidates.append(resolved)

    for candidate in deduped_candidates:
        if candidate.exists():
            return candidate

    if deduped_candidates:
        return deduped_candidates[0]

    raise ValueError("Uploaded file path is outside the configured upload directory")


def _extract_pdf(path: Path) -> str:
    with fitz.open(path) as document:
        return "\n\n".join(page.get_text("text") for page in document)


def _extract_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("PPTX extraction requires python-pptx to be installed") from exc

    presentation = Presentation(str(path))
    parts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"Slide {slide_index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_parts.append(shape.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                slide_parts.append(f"Speaker notes: {notes}")
        parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


def _extract_text(path: Path, file_kind: str, mime_type: str) -> str:
    if file_kind == "pdf" or mime_type == "application/pdf":
        return _extract_pdf(path)
    if file_kind == "txt" or mime_type == "text/plain":
        return _extract_txt(path)
    if file_kind == "pptx":
        return _extract_pptx(path)
    raise ValueError(f"Unsupported library file kind: {file_kind}")


async def delete_library_file_chunks(db: AsyncSession, file_id: str) -> dict[str, Any]:
    result = await db.execute(
        sa_text(
            """
            DELETE FROM content_chunks
            WHERE library_file_id = :fileId
               OR (source_type = 'library_file' AND source_id = :fileId)
            """
        ),
        {"fileId": file_id},
    )
    await db.commit()
    return {"fileId": file_id, "chunksDeleted": result.rowcount or 0}


async def index_library_file(db: AsyncSession, file_id: str) -> dict[str, Any]:
    row_result = await db.execute(
        sa_text(
            """
            SELECT
              id,
              file_path,
              original_name,
              mime_type,
              size_bytes,
              teacher_id,
              class_id,
              scope,
              ai_enabled,
              subject_key,
              grade_level,
              teacher_visible,
              file_kind,
              content_hash
            FROM uploaded_files
            WHERE id = :fileId
              AND deleted_at IS NULL
              AND (
                scope = 'general'
                OR (
                  scope = 'private'
                  AND ai_enabled = true
                )
              )
            """
        ),
        {"fileId": file_id},
    )
    file_row = row_result.mappings().first()
    if not file_row:
        raise ValueError("Library file not found")
    if not file_row["subject_key"] or not file_row["grade_level"]:
        raise ValueError("Library file is missing subject_key or grade_level")

    await db.execute(
        sa_text(
            """
            UPDATE uploaded_files
            SET index_status = 'processing',
                index_error = NULL
            WHERE id = :fileId
            """
        ),
        {"fileId": file_id},
    )
    await db.commit()

    try:
        materialized_path = await materialize_backend_upload(str(file_row["file_path"]))
        if not materialized_path:
            raise FileNotFoundError(
                f"File could not be materialized from backend storage: {file_row['file_path']}"
            )
        resolved_path = Path(materialized_path).resolve()
        if not resolved_path.exists():
            raise FileNotFoundError(f"File does not exist on disk: {resolved_path}")

        raw_text = _extract_text(
            resolved_path,
            str(file_row["file_kind"] or ""),
            str(file_row["mime_type"] or ""),
        )
        sanitized = sanitize_extracted_text(raw_text)
        chunks = chunk_text_for_indexing(sanitized.cleaned_text)

        await delete_library_file_chunks(db, file_id)

        if not chunks:
            raise ValueError("No indexable text was extracted from the library file")

        embeddings = await embed_texts(chunks)
        created = 0
        for chunk_order, (chunk_text, embedding) in enumerate(zip(chunks, embeddings)):
            content_hash = hashlib.sha256(
                f"library_file:{file_id}:{chunk_order}:{chunk_text}".encode("utf-8")
            ).hexdigest()
            metadata = {
                "documentId": f"library:{file_id}:chunk:{chunk_order}",
                "libraryFileId": file_id,
                "title": file_row["original_name"],
                "originalName": file_row["original_name"],
                "sourceReference": f"library:{file_id} | chunk:{chunk_order}",
                "sourceType": "library_file",
                "blockType": "library_file",
                "teacherId": str(file_row["teacher_id"]) if file_row["teacher_id"] else None,
                "classId": str(file_row["class_id"]) if file_row["class_id"] else None,
                "scope": file_row["scope"],
                "aiEnabled": bool(file_row["ai_enabled"]),
                "subjectKey": file_row["subject_key"],
                "gradeLevel": file_row["grade_level"],
                "teacherVisible": bool(file_row["teacher_visible"]),
                "isPublished": True,
                "fileKind": file_row["file_kind"],
                "mimeType": file_row["mime_type"],
                "sizeBytes": int(file_row["size_bytes"] or 0),
                "sanitizationWarnings": sanitized.warnings,
            }

            insert_result = await db.execute(
                sa_text(
                    """
                    INSERT INTO content_chunks (
                      source_type,
                      source_id,
                      class_id,
                      library_file_id,
                      subject_key,
                      grade_level,
                      chunk_text,
                      chunk_order,
                      token_count,
                      content_hash,
                      metadata_json
                    )
                    VALUES (
                      'library_file',
                      :sourceId,
                      NULL,
                      :libraryFileId,
                      :subjectKey,
                      :gradeLevel,
                      :chunkText,
                      :chunkOrder,
                      :tokenCount,
                      :contentHash,
                      :metadataJson
                    )
                    RETURNING id
                    """
                ).bindparams(bindparam("metadataJson", type_=postgresql.JSONB)),
                {
                    "sourceId": file_id,
                    "libraryFileId": file_id,
                    "subjectKey": file_row["subject_key"],
                    "gradeLevel": file_row["grade_level"],
                    "chunkText": chunk_text,
                    "chunkOrder": chunk_order,
                    "tokenCount": estimate_token_count(chunk_text),
                    "contentHash": content_hash,
                    "metadataJson": metadata,
                },
            )
            chunk_id = insert_result.scalar_one()
            await db.execute(
                sa_text(
                    """
                    INSERT INTO content_chunk_embeddings (
                      chunk_id,
                      embedding,
                      embedding_model,
                      embedded_at
                    )
                    VALUES (
                      :chunkId,
                      CAST(:embedding AS vector),
                      :embeddingModel,
                      NOW()
                    )
                    """
                ),
                {
                    "chunkId": chunk_id,
                    "embedding": embedding_to_vector_literal(embedding),
                    "embeddingModel": "ollama:" + ollama_client.get_embedding_model_name(),
                },
            )
            created += 1

        await db.execute(
            sa_text(
                """
                UPDATE uploaded_files
                SET index_status = 'completed',
                    index_error = NULL,
                    indexed_at = NOW()
                WHERE id = :fileId
                """
            ),
            {"fileId": file_id},
        )
        await db.commit()
        logger.info("[library-index] Indexed %s with %d chunk(s)", file_id, created)
        return {
            "fileId": file_id,
            "chunksIndexed": created,
            "sanitizationWarnings": sanitized.warnings,
        }
    except Exception as exc:
        await db.execute(
            sa_text(
                """
                UPDATE uploaded_files
                SET index_status = 'failed',
                    index_error = :error
                WHERE id = :fileId
                """
            ),
            {"fileId": file_id, "error": str(exc)[:1000]},
        )
        await db.commit()
        raise


async def backfill_library_files(db: AsyncSession) -> dict[str, Any]:
    rows = await db.execute(
        sa_text(
            """
            SELECT id
            FROM uploaded_files
            WHERE deleted_at IS NULL
              AND (
                scope = 'general'
                OR (
                  scope = 'private'
                  AND ai_enabled = true
                )
              )
              AND subject_key IS NOT NULL
              AND grade_level IS NOT NULL
            ORDER BY uploaded_at ASC
            """
        )
    )
    file_ids = [str(row["id"]) for row in rows.mappings()]
    results = []
    for file_id in file_ids:
        results.append(await index_library_file(db, file_id))
    return {
        "filesProcessed": len(results),
        "results": results,
    }
